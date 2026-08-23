# SPDX-License-Identifier: FSL-1.1-ALv2
# Copyright 2026 Alexander Atkins
# Licensed under the Functional Source License, Version 1.1, ALv2 Future License.
# Each version becomes Apache-2.0 two years after its release. See LICENSE.

"""
Document format extractors for non-HTML content types.

Supports: PDF, DOCX, XLSX, PPTX
Each extractor returns (title, body) strings.
"""

from __future__ import annotations

import io
import logging

log = logging.getLogger(__name__)


def extract_pdf(data: bytes) -> tuple[str, str]:
    """Extract text from a PDF using pdfplumber. Returns (title, body)."""
    import pdfplumber

    pages = []
    title = ""
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        # Try to get title from metadata
        meta = pdf.metadata or {}
        title = meta.get("Title") or meta.get("Subject") or ""
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                pages.append(text)
            if i > 100:  # cap at 100 pages — avoid monster docs
                log.debug("PDF truncated at 100 pages")
                break
    body = "\n\n".join(pages)
    # If no metadata title, use first non-empty line
    if not title and body:
        title = body.splitlines()[0][:120].strip()
    return title, body


def extract_docx(data: bytes) -> tuple[str, str]:
    """Extract text from a DOCX file. Returns (title, body)."""
    from docx import Document

    doc = Document(io.BytesIO(data))
    # First heading or first paragraph as title
    title = ""
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if not title and para.style.name.startswith("Heading"):
            title = text
        paragraphs.append(text)
    body = "\n\n".join(paragraphs)
    if not title and paragraphs:
        title = paragraphs[0][:120]
    return title, body


def extract_xlsx(data: bytes) -> tuple[str, str]:
    """Extract text from an XLSX spreadsheet. Returns (title, body)."""
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    title = wb.properties.title or ""
    sheets = []
    for sheet_name in wb.sheetnames[:10]:  # cap at 10 sheets
        ws = wb[sheet_name]
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i > 500:  # cap at 500 rows per sheet
                break
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append("\t".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    body = "\n\n".join(sheets)
    if not title:
        title = wb.sheetnames[0] if wb.sheetnames else "Spreadsheet"
    return title, body


def extract_pptx(data: bytes) -> tuple[str, str]:
    """Extract text from a PPTX presentation. Returns (title, body)."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    title = ""
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
        if texts:
            # First shape on slide 1 is usually the presentation title
            if i == 0 and not title:
                title = texts[0][:120]
            slides.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
    body = "\n\n".join(slides)
    return title, body


# MIME type / extension → extractor mapping
EXTRACTORS: dict[str, callable] = {
    "application/pdf": extract_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": extract_xlsx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_pptx,
    "application/msword": extract_docx,  # legacy .doc — best effort
}

EXTENSION_MAP: dict[str, str] = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc": "application/msword",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def get_extractor(content_type: str, url: str) -> callable | None:
    """Return the right extractor fn or None if not a supported binary format."""
    # Normalise content_type (strip charset etc)
    ct = content_type.split(";")[0].strip().lower()
    if ct in EXTRACTORS:
        return EXTRACTORS[ct]
    # Fallback: check URL extension
    import os
    from urllib.parse import urlparse

    ext = os.path.splitext(urlparse(url).path)[1].lower()
    mime = EXTENSION_MAP.get(ext)
    if mime:
        return EXTRACTORS.get(mime)
    return None
